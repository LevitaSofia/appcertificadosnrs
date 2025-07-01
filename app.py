import os
import re
import subprocess
import comtypes.client
from flask import Flask, render_template, request, redirect, url_for, flash, send_file, jsonify, Response, session
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from datetime import datetime
import shutil
from pptx import Presentation
from werkzeug.utils import secure_filename
import sqlite3
import mimetypes
from werkzeug.exceptions import NotFound

app = Flask(__name__)
app.config['SECRET_KEY'] = 'sua-chave-secreta-aqui'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///certificados.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = 'static/imagens'
app.config['VIDEOS_FOLDER'] = 'static/videos'
app.config['MAX_CONTENT_LENGTH'] = 500 * \
    1024 * 1024  # 500MB máximo por arquivo

db = SQLAlchemy(app)

# Configurar Flask-Login
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = 'Você precisa fazer login para acessar esta página.'
login_manager.login_message_category = 'warning'

# Filtros Jinja2 personalizados


@app.template_filter('basename')
def basename_filter(filepath):
    """Extrai o nome do arquivo de um caminho completo"""
    if not filepath:
        return ''
    return os.path.basename(filepath)

# Modelos do Banco de Dados


class Funcionario(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    nome = db.Column(db.String(100), nullable=False)
    cpf = db.Column(db.String(14), unique=True, nullable=False)
    rg = db.Column(db.String(20), nullable=True)
    funcao = db.Column(db.String(100), nullable=True, default='Não informado')
    data_admissao = db.Column(db.Date, nullable=True)
    data_nascimento = db.Column(db.Date, nullable=True)
    telefone = db.Column(db.String(20), nullable=True)
    email = db.Column(db.String(100), nullable=True)
    senha_hash = db.Column(db.String(200), nullable=True)
    foto = db.Column(db.String(200))
    ativo = db.Column(db.Boolean, default=True, nullable=False)
    admin = db.Column(db.Boolean, default=False, nullable=False)
    primeiro_login = db.Column(db.Boolean, default=True, nullable=False)
    data_ultimo_login = db.Column(db.DateTime, nullable=True)

    def set_password(self, password):
        """Define a senha do funcionário"""
        self.senha_hash = generate_password_hash(password)

    def check_password(self, password):
        """Verifica a senha do funcionário"""
        if not self.senha_hash:
            return False
        return check_password_hash(self.senha_hash, password)

    def is_active(self):
        """Retorna True se o funcionário está ativo"""
        return self.ativo

    def get_id(self):
        """Retorna o ID do usuário como string"""
        return str(self.id)

    @property
    def cpf_formatado(self):
        """Retorna CPF formatado"""
        if len(self.cpf) == 11:
            return f"{self.cpf[:3]}.{self.cpf[3:6]}.{self.cpf[6:9]}-{self.cpf[9:]}"
        return self.cpf

    certificados = db.relationship(
        'CertificadoGerado', backref='funcionario', lazy=True)
    progresso = db.relationship(
        'ProgressoTreinamento', backref='funcionario', lazy=True)
    resultados_avaliacao = db.relationship(
        'ResultadoAvaliacao', backref='funcionario', lazy=True)


class CertificadoGerado(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    funcionario_id = db.Column(db.Integer, db.ForeignKey(
        'funcionario.id'), nullable=False)
    tipo_nr = db.Column(db.String(10), nullable=False)
    tipo_treinamento = db.Column(db.String(50), nullable=False)
    data_emissao = db.Column(db.Date, nullable=False)
    caminho_arquivo = db.Column(db.String(300), nullable=False)


class ModeloNR(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    tipo_nr = db.Column(db.String(10), unique=True, nullable=False)
    caminho_modelo_pptx = db.Column(db.String(300), nullable=False)
    descricao = db.Column(db.String(200))


class Cargo(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    nome = db.Column(db.String(100), unique=True, nullable=False)
    descricao = db.Column(db.String(200), nullable=True)
    ativo = db.Column(db.Boolean, default=True, nullable=False)
    data_criacao = db.Column(
        db.Date, default=datetime.now().date, nullable=False)


class Treinamento(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    titulo = db.Column(db.String(200), nullable=False)
    descricao = db.Column(db.Text, nullable=True)
    tipo_nr = db.Column(db.String(10), nullable=False)
    duracao_minutos = db.Column(db.Integer, nullable=True)
    arquivo_video = db.Column(db.String(500), nullable=False)
    thumbnail = db.Column(db.String(500), nullable=True)
    ativo = db.Column(db.Boolean, default=True, nullable=False)
    obrigatorio = db.Column(db.Boolean, default=False, nullable=False)
    data_criacao = db.Column(db.DateTime, default=datetime.now, nullable=False)
    data_atualizacao = db.Column(
        db.DateTime, default=datetime.now, onupdate=datetime.now, nullable=False)

    # Relacionamentos
    progresso = db.relationship(
        'ProgressoTreinamento', backref='treinamento', lazy=True, cascade='all, delete-orphan')


class ProgressoTreinamento(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    funcionario_id = db.Column(db.Integer, db.ForeignKey(
        'funcionario.id'), nullable=False)
    treinamento_id = db.Column(db.Integer, db.ForeignKey(
        'treinamento.id'), nullable=False)
    progresso_percent = db.Column(
        db.Float, default=0.0, nullable=False)  # 0-100
    tempo_assistido_segundos = db.Column(db.Integer, default=0, nullable=False)
    concluido = db.Column(db.Boolean, default=False, nullable=False)
    data_inicio = db.Column(db.DateTime, default=datetime.now, nullable=False)
    data_conclusao = db.Column(db.DateTime, nullable=True)
    data_ultima_visualizacao = db.Column(
        db.DateTime, default=datetime.now, nullable=False)

    # Constraint única: um funcionário só pode ter um progresso por treinamento
    __table_args__ = (db.UniqueConstraint(
        'funcionario_id', 'treinamento_id', name='unique_funcionario_treinamento'),)


class Avaliacao(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    treinamento_id = db.Column(db.Integer, db.ForeignKey('treinamento.id'), nullable=False)
    titulo = db.Column(db.String(200), nullable=False)
    descricao = db.Column(db.Text, nullable=True)
    nota_minima_aprovacao = db.Column(db.Float, default=7.0, nullable=False)
    ativo = db.Column(db.Boolean, default=True, nullable=False)
    data_criacao = db.Column(db.DateTime, default=datetime.now, nullable=False)
    
    # Relacionamentos
    perguntas = db.relationship('PerguntaAvaliacao', backref='avaliacao', lazy=True, cascade='all, delete-orphan')
    resultados = db.relationship('ResultadoAvaliacao', backref='avaliacao', lazy=True, cascade='all, delete-orphan')


class PerguntaAvaliacao(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    avaliacao_id = db.Column(db.Integer, db.ForeignKey('avaliacao.id'), nullable=False)
    texto_pergunta = db.Column(db.Text, nullable=False)
    tipo_pergunta = db.Column(db.String(20), default='multipla_escolha', nullable=False)  # multipla_escolha, verdadeiro_falso
    ordem = db.Column(db.Integer, nullable=False)
    pontos = db.Column(db.Float, default=1.0, nullable=False)
    
    # Relacionamentos
    opcoes = db.relationship('OpcaoResposta', backref='pergunta', lazy=True, cascade='all, delete-orphan')


class OpcaoResposta(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    pergunta_id = db.Column(db.Integer, db.ForeignKey('pergunta_avaliacao.id'), nullable=False)
    texto_opcao = db.Column(db.Text, nullable=False)
    correta = db.Column(db.Boolean, default=False, nullable=False)
    ordem = db.Column(db.Integer, nullable=False)


class ResultadoAvaliacao(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    funcionario_id = db.Column(db.Integer, db.ForeignKey('funcionario.id'), nullable=False)
    avaliacao_id = db.Column(db.Integer, db.ForeignKey('avaliacao.id'), nullable=False)
    nota_obtida = db.Column(db.Float, nullable=False)
    aprovado = db.Column(db.Boolean, nullable=False)
    data_realizacao = db.Column(db.DateTime, default=datetime.now, nullable=False)
    tempo_realizacao_segundos = db.Column(db.Integer, nullable=True)
    
    # Constraint única: um funcionário só pode ter um resultado por avaliação
    __table_args__ = (db.UniqueConstraint('funcionario_id', 'avaliacao_id', name='unique_funcionario_avaliacao'),)


class RespostaFuncionario(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    resultado_avaliacao_id = db.Column(db.Integer, db.ForeignKey('resultado_avaliacao.id'), nullable=False)
    pergunta_id = db.Column(db.Integer, db.ForeignKey('pergunta_avaliacao.id'), nullable=False)
    opcao_escolhida_id = db.Column(db.Integer, db.ForeignKey('opcao_resposta.id'), nullable=True)
    correta = db.Column(db.Boolean, nullable=False)
    
    # Relacionamentos
    resultado = db.relationship('ResultadoAvaliacao', backref='respostas')
    pergunta = db.relationship('PerguntaAvaliacao')
    opcao_escolhida = db.relationship('OpcaoResposta')

# Função para validar CPF


def validar_cpf(cpf):
    cpf = re.sub(r'[^0-9]', '', cpf)
    if len(cpf) != 11:
        return False

    if cpf == cpf[0] * 11:
        return False

    def calcular_digito(cpf, peso):
        soma = sum(int(digito) * peso for digito,
                   peso in zip(cpf, range(peso, 1, -1)))
        resto = soma % 11
        return 0 if resto < 2 else 11 - resto

    primeiro_digito = calcular_digito(cpf, 10)
    segundo_digito = calcular_digito(cpf, 11)

    return cpf[-2:] == f"{primeiro_digito}{segundo_digito}"

# Função para gerar certificado em PDF


# Função para converter PowerPoint para PDF
def converter_pptx_para_pdf(caminho_pptx, caminho_pdf):
    try:
        print(f"Tentando converter: {caminho_pptx} -> {caminho_pdf}")

        # Método 1: Tentar usar COM (PowerPoint)
        try:
            import comtypes.client

            # Criar aplicação PowerPoint
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = False  # Executar em background

            # Abrir apresentação
            presentation = powerpoint.Presentations.Open(
                os.path.abspath(caminho_pptx))

            # Salvar como PDF
            presentation.SaveAs(os.path.abspath(
                caminho_pdf), 32)  # 32 = ppSaveAsPDF

            # Fechar apresentação e PowerPoint
            presentation.Close()
            powerpoint.Quit()

            print("Conversão COM bem-sucedida!")
            return True

        except Exception as com_error:
            print(f"Erro COM: {str(com_error)}")

            # Método 2: Tentar usar win32com como alternativa
            try:
                import win32com.client

                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = 0  # Não mostrar PowerPoint

                presentation = powerpoint.Presentations.Open(
                    os.path.abspath(caminho_pptx))
                presentation.SaveAs(os.path.abspath(caminho_pdf), 32)
                presentation.Close()
                powerpoint.Quit()

                print("Conversão win32com bem-sucedida!")
                return True

            except Exception as win32_error:
                print(f"Erro win32com: {str(win32_error)}")

                # Método 3: Usar subprocess como último recurso
                try:
                    import subprocess

                    # Comando PowerShell para converter
                    powershell_cmd = f'''
                    $powerpoint = New-Object -ComObject PowerPoint.Application
                    $powerpoint.Visible = $false
                    $presentation = $powerpoint.Presentations.Open("{os.path.abspath(caminho_pptx)}")
                    $presentation.SaveAs("{os.path.abspath(caminho_pdf)}", 32)
                    $presentation.Close()
                    $powerpoint.Quit()
                    '''

                    result = subprocess.run(
                        ["powershell", "-Command", powershell_cmd],
                        capture_output=True,
                        text=True,
                        timeout=60
                    )

                    if result.returncode == 0:
                        print("Conversão PowerShell bem-sucedida!")
                        return True
                    else:
                        print(f"Erro PowerShell: {result.stderr}")

                except Exception as ps_error:
                    print(f"Erro PowerShell: {str(ps_error)}")

        return False

    except Exception as e:
        print(f"Erro geral na conversão: {str(e)}")
        return False

# Função para gerar certificado usando PowerPoint como base


def gerar_certificado(funcionario, tipo_nr, tipo_treinamento, data_emissao):
    try:
        # Buscar modelo
        modelo = ModeloNR.query.filter_by(tipo_nr=tipo_nr).first()
        if not modelo:
            return False, "Modelo não encontrado para este tipo de NR"

        # Verificar se arquivo modelo existe
        if not os.path.exists(modelo.caminho_modelo_pptx):
            return False, f"Arquivo do modelo não encontrado: {modelo.caminho_modelo_pptx}"

        # Carregar apresentação
        prs = Presentation(modelo.caminho_modelo_pptx)

        # Criar pasta do funcionário se não existir
        nome_funcionario_limpo = re.sub(r'[<>:"/\\|?*]', '_', funcionario.nome)
        pasta_funcionario = os.path.join(
            'certificados', nome_funcionario_limpo)
        os.makedirs(pasta_funcionario, exist_ok=True)

        # Preparar substituições
        data_formatada = data_emissao.strftime('%d/%m/%Y')
        data_admissao_formatada = funcionario.data_admissao.strftime(
            '%d/%m/%Y') if funcionario.data_admissao else 'Não informado'

        substituicoes = {
            '{{NOME}}': funcionario.nome,
            '{{CPF}}': funcionario.cpf,
            '{{RG}}': funcionario.rg if funcionario.rg else 'Não informado',
            '{{CARGO}}': funcionario.funcao if funcionario.funcao else 'Não informado',
            '{{FUNCAO}}': funcionario.funcao if funcionario.funcao else 'Não informado',
            '{{TIPO_TREINAMENTO}}': tipo_treinamento,
            '{{DATA}}': data_formatada,
            '{{DATA_EMISSAO}}': data_formatada,
            '{{DATA_ADMISSAO}}': data_admissao_formatada,
            '{{NR}}': tipo_nr,
            '{{TIPO_NR}}': tipo_nr,
            '{{DESCRICAO_NR}}': modelo.descricao
        }

        # Substituir texto nos slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for find_text, replace_text in substituicoes.items():
                        if find_text in shape.text:
                            shape.text = shape.text.replace(
                                find_text, replace_text)

                # Verificar text_frame
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for find_text, replace_text in substituicoes.items():
                                if find_text in run.text:
                                    run.text = run.text.replace(
                                        find_text, replace_text)

        # Salvar PowerPoint temporário
        data_formatada_arquivo = data_emissao.strftime('%Y-%m-%d')
        nome_arquivo_base = f"{nome_funcionario_limpo}_{data_formatada_arquivo}_{tipo_nr}_{tipo_treinamento}"
        caminho_pptx_temp = os.path.join(
            pasta_funcionario, f"{nome_arquivo_base}.pptx")
        caminho_pdf_final = os.path.join(
            pasta_funcionario, f"{nome_arquivo_base}.pdf")

        # Salvar PowerPoint preenchido
        print(f"Salvando PowerPoint em: {caminho_pptx_temp}")
        prs.save(caminho_pptx_temp)

        # Converter para PDF
        print("Iniciando conversão para PDF...")
        if converter_pptx_para_pdf(caminho_pptx_temp, caminho_pdf_final):
            print("✅ Conversão para PDF bem-sucedida!")
            # Remover arquivo PowerPoint temporário
            try:
                os.remove(caminho_pptx_temp)
                print("Arquivo PowerPoint temporário removido")
            except:
                print("Não foi possível remover o arquivo PowerPoint temporário")
                pass  # Se não conseguir remover, não é crítico

            caminho_final = caminho_pdf_final
            tipo_arquivo = "PDF"
        else:
            print("⚠️ Não foi possível converter para PDF, mantendo arquivo PowerPoint")
            caminho_final = caminho_pptx_temp
            tipo_arquivo = "PowerPoint"

        # Registrar no banco
        certificado = CertificadoGerado(
            funcionario_id=funcionario.id,
            tipo_nr=tipo_nr,
            tipo_treinamento=tipo_treinamento,
            data_emissao=data_emissao,
            caminho_arquivo=caminho_final
        )
        db.session.add(certificado)
        db.session.commit()

        print(f"✅ Certificado gerado com sucesso em: {caminho_final}")
        return True, caminho_final

    except Exception as e:
        print(f"❌ Erro ao gerar certificado: {str(e)}")
        return False, f"Erro ao gerar certificado: {str(e)}"

# Rotas


@app.route('/')
def index():
    funcionarios = Funcionario.query.all()
    modelos = ModeloNR.query.all()
    return render_template('index.html', funcionarios=funcionarios, modelos=modelos)


@app.route('/funcionarios')
def funcionarios():
    funcionarios = Funcionario.query.all()
    return render_template('funcionarios.html', funcionarios=funcionarios)


@app.route('/cadastrar_funcionario', methods=['GET', 'POST'])
def cadastrar_funcionario():
    if request.method == 'POST':
        nome = request.form['nome']
        cpf = request.form['cpf']
        rg = request.form.get('rg', '')
        funcao = request.form.get('funcao', 'Não informado')
        data_admissao_str = request.form.get('data_admissao')
        data_nascimento_str = request.form.get('data_nascimento')
        telefone = request.form.get('telefone', '')
        email = request.form.get('email', '')
        senha = request.form.get('senha', '')

        # Converter datas
        data_admissao = None
        if data_admissao_str:
            data_admissao = datetime.strptime(
                data_admissao_str, '%Y-%m-%d').date()

        data_nascimento = None
        if data_nascimento_str:
            data_nascimento = datetime.strptime(
                data_nascimento_str, '%Y-%m-%d').date()

        # Validar CPF
        if not validar_cpf(cpf):
            flash('CPF inválido!', 'error')
            return render_template('cadastrar_funcionario.html')

        # Verificar se CPF já existe
        if Funcionario.query.filter_by(cpf=cpf).first():
            flash('CPF já cadastrado!', 'error')
            return render_template('cadastrar_funcionario.html')

        # Salvar funcionário
        funcionario = Funcionario(
            nome=nome,
            cpf=cpf,
            rg=rg,
            funcao=funcao,
            data_admissao=data_admissao,
            data_nascimento=data_nascimento,
            telefone=telefone,
            email=email
        )
        
        # Define a senha se fornecida
        if senha:
            funcionario.set_password(senha)

        db.session.add(funcionario)
        db.session.commit()

        flash('Funcionário cadastrado com sucesso!', 'success')
        return redirect(url_for('funcionarios'))

    cargos = Cargo.query.filter_by(ativo=True).order_by(Cargo.nome).all()
    return render_template('cadastrar_funcionario.html', cargos=cargos)


@app.route('/editar_funcionario/<int:funcionario_id>', methods=['GET', 'POST'])
def editar_funcionario(funcionario_id):
    funcionario = Funcionario.query.get_or_404(funcionario_id)

    if request.method == 'POST':
        funcionario.nome = request.form['nome']
        funcionario.rg = request.form.get('rg', '')
        funcionario.funcao = request.form.get('funcao', 'Não informado')
        funcionario.telefone = request.form.get('telefone', '')
        funcionario.email = request.form.get('email', '')
        funcionario.senha = request.form.get('senha', '')

        # Converter datas
        data_admissao_str = request.form.get('data_admissao')
        if data_admissao_str:
            funcionario.data_admissao = datetime.strptime(
                data_admissao_str, '%Y-%m-%d').date()

        data_nascimento_str = request.form.get('data_nascimento')
        if data_nascimento_str:
            funcionario.data_nascimento = datetime.strptime(
                data_nascimento_str, '%Y-%m-%d').date()

        # Validar CPF apenas se foi alterado
        cpf_novo = request.form['cpf']
        if cpf_novo != funcionario.cpf:
            if not validar_cpf(cpf_novo):
                flash('CPF inválido!', 'error')
                cargos = Cargo.query.filter_by(
                    ativo=True).order_by(Cargo.nome).all()
                return render_template('editar_funcionario.html', funcionario=funcionario, cargos=cargos)

            # Verificar se CPF já existe
            if Funcionario.query.filter_by(cpf=cpf_novo).first():
                flash('CPF já cadastrado por outro funcionário!', 'error')
                return render_template('editar_funcionario.html', funcionario=funcionario)

            funcionario.cpf = cpf_novo

        db.session.commit()
        flash('Funcionário atualizado com sucesso!', 'success')
        return redirect(url_for('funcionarios'))

    cargos = Cargo.query.filter_by(ativo=True).order_by(Cargo.nome).all()
    return render_template('editar_funcionario.html', funcionario=funcionario, cargos=cargos)


@app.route('/excluir_funcionario/<int:funcionario_id>', methods=['POST'])
def excluir_funcionario(funcionario_id):
    funcionario = Funcionario.query.get_or_404(funcionario_id)

    # Verificar se tem certificados
    if funcionario.certificados:
        flash(
            f'Não é possível excluir {funcionario.nome}. Funcionário possui certificados gerados.', 'error')
        return redirect(url_for('funcionarios'))

    nome_funcionario = funcionario.nome
    db.session.delete(funcionario)
    db.session.commit()

    flash(f'Funcionário {nome_funcionario} excluído com sucesso!', 'success')
    return redirect(url_for('funcionarios'))


@app.route('/gerar_unitario')
def gerar_unitario():
    funcionarios = Funcionario.query.all()
    modelos = ModeloNR.query.all()
    data_hoje = datetime.now().strftime('%Y-%m-%d')
    return render_template('gerar_unitario.html', funcionarios=funcionarios, modelos=modelos, data_hoje=data_hoje)


@app.route('/processar_certificado_unitario', methods=['POST'])
def processar_certificado_unitario():
    funcionario_id = request.form['funcionario_id']
    tipo_nr = request.form['tipo_nr']
    tipo_treinamento = request.form['tipo_treinamento']
    data_emissao = datetime.strptime(
        request.form['data_emissao'], '%Y-%m-%d').date()

    funcionario = Funcionario.query.get_or_404(funcionario_id)

    sucesso, resultado = gerar_certificado(
        funcionario, tipo_nr, tipo_treinamento, data_emissao)

    if sucesso:
        flash('Certificado gerado com sucesso!', 'success')
        return send_file(resultado, as_attachment=True)
    else:
        flash(f'Erro: {resultado}', 'error')
        return redirect(url_for('gerar_unitario'))


@app.route('/gerar_lote')
def gerar_lote():
    funcionarios = Funcionario.query.all()
    modelos = ModeloNR.query.all()
    data_hoje = datetime.now().strftime('%Y-%m-%d')
    return render_template('gerar_lote.html', funcionarios=funcionarios, modelos=modelos, data_hoje=data_hoje)


@app.route('/processar_certificados_lote', methods=['POST'])
def processar_certificados_lote():
    funcionarios_ids = request.form.getlist('funcionarios_ids')
    tipo_nr = request.form['tipo_nr']
    tipo_treinamento = request.form['tipo_treinamento']
    data_emissao = datetime.strptime(
        request.form['data_emissao'], '%Y-%m-%d').date()

    resultados = []

    for funcionario_id in funcionarios_ids:
        funcionario = Funcionario.query.get(funcionario_id)
        if funcionario:
            sucesso, resultado = gerar_certificado(
                funcionario, tipo_nr, tipo_treinamento, data_emissao)
            resultados.append({
                'funcionario': funcionario.nome,
                'sucesso': sucesso,
                'resultado': resultado
            })

    flash(f'Processados {len(resultados)} certificados', 'info')
    return render_template('resultado_lote.html', resultados=resultados)


@app.route('/relatorios')
def relatorios():
    # Estatísticas gerais
    total_funcionarios = Funcionario.query.count()
    total_certificados = CertificadoGerado.query.count()

    # Certificados por NR
    certificados_por_nr = db.session.query(
        CertificadoGerado.tipo_nr,
        db.func.count(CertificadoGerado.id).label('total')
    ).group_by(CertificadoGerado.tipo_nr).all()

    # Últimos certificados
    ultimos_certificados = db.session.query(CertificadoGerado, Funcionario).join(
        Funcionario, CertificadoGerado.funcionario_id == Funcionario.id
    ).order_by(CertificadoGerado.data_emissao.desc()).limit(10).all()

    return render_template('relatorios.html',
                           total_funcionarios=total_funcionarios,
                           total_certificados=total_certificados,
                           certificados_por_nr=certificados_por_nr,
                           ultimos_certificados=ultimos_certificados)


@app.route('/configuracoes')
def configuracoes():
    modelos = ModeloNR.query.all()
    cargos = Cargo.query.filter_by(ativo=True).all()
    return render_template('configuracoes.html', modelos=modelos, cargos=cargos)


@app.route('/cadastrar_modelo', methods=['POST'])
def cadastrar_modelo():
    tipo_nr = request.form['tipo_nr']
    descricao = request.form['descricao']

    # Verificar se já existe
    if ModeloNR.query.filter_by(tipo_nr=tipo_nr).first():
        flash('Modelo para esta NR já existe!', 'error')
        return redirect(url_for('configuracoes'))

    # Caminho do modelo
    caminho_modelo = f"modelos_nr/{tipo_nr}_modelo.pptx"

    modelo = ModeloNR(
        tipo_nr=tipo_nr,
        caminho_modelo_pptx=caminho_modelo,
        descricao=descricao
    )

    db.session.add(modelo)
    db.session.commit()

    flash('Modelo cadastrado! Lembre-se de colocar o arquivo PowerPoint na pasta modelos_nr/', 'success')
    return redirect(url_for('configuracoes'))


@app.route('/cadastrar_cargo', methods=['POST'])
def cadastrar_cargo():
    nome = request.form['nome'].strip()
    descricao = request.form.get('descricao', '').strip()

    if not nome:
        flash('Nome do cargo é obrigatório!', 'error')
        return redirect(url_for('configuracoes'))

    # Verificar se já existe
    if Cargo.query.filter_by(nome=nome, ativo=True).first():
        flash('Cargo já existe!', 'error')
        return redirect(url_for('configuracoes'))

    cargo = Cargo(
        nome=nome,
        descricao=descricao if descricao else None
    )

    db.session.add(cargo)
    db.session.commit()

    flash(f'Cargo "{nome}" cadastrado com sucesso!', 'success')
    return redirect(url_for('configuracoes'))


@app.route('/editar_cargo/<int:cargo_id>', methods=['POST'])
def editar_cargo(cargo_id):
    cargo = Cargo.query.get_or_404(cargo_id)

    nome = request.form['nome'].strip()
    descricao = request.form.get('descricao', '').strip()

    if not nome:
        flash('Nome do cargo é obrigatório!', 'error')
        return redirect(url_for('configuracoes'))

    # Verificar se já existe outro cargo com mesmo nome
    cargo_existente = Cargo.query.filter_by(nome=nome, ativo=True).first()
    if cargo_existente and cargo_existente.id != cargo_id:
        flash('Já existe outro cargo com este nome!', 'error')
        return redirect(url_for('configuracoes'))

    cargo.nome = nome
    cargo.descricao = descricao if descricao else None

    db.session.commit()

    flash(f'Cargo "{nome}" atualizado com sucesso!', 'success')
    return redirect(url_for('configuracoes'))


@app.route('/excluir_cargo/<int:cargo_id>', methods=['POST'])
def excluir_cargo(cargo_id):
    cargo = Cargo.query.get_or_404(cargo_id)

    # Verificar se há funcionários usando este cargo
    funcionarios_com_cargo = Funcionario.query.filter_by(
        funcao=cargo.nome).count()

    if funcionarios_com_cargo > 0:
        flash(
            f'Não é possível excluir o cargo "{cargo.nome}". Existem {funcionarios_com_cargo} funcionário(s) com este cargo.', 'error')
        return redirect(url_for('configuracoes'))

    # Desativar ao invés de excluir (soft delete)
    cargo.ativo = False
    db.session.commit()

    flash(f'Cargo "{cargo.nome}" removido com sucesso!', 'success')
    return redirect(url_for('configuracoes'))


@app.route('/api/cargos')
def api_cargos():
    """API para listar cargos ativos - usado nos formulários"""
    cargos = Cargo.query.filter_by(ativo=True).order_by(Cargo.nome).all()
    return jsonify([{'id': c.id, 'nome': c.nome, 'descricao': c.descricao} for c in cargos])


# ==================== CENTRAL DE TREINAMENTOS ====================

@app.route('/treinamentos', endpoint='treinamentos')
@login_required
def modulos_treinamento():
    """Página que exibe os módulos de NR disponíveis."""
    modulos = db.session.query(Treinamento.tipo_nr).distinct().order_by(
        Treinamento.tipo_nr).all()
    # Extrai os nomes dos módulos da tupla
    modulos_lista = [m[0] for m in modulos]
    return render_template('modulos_treinamento.html', modulos=modulos_lista)


@app.route('/treinamentos/<string:tipo_nr>')
def treinamentos_por_modulo(tipo_nr):
    """Página que exibe os treinamentos de um módulo de NR específico."""
    treinamentos = Treinamento.query.filter_by(
        tipo_nr=tipo_nr, ativo=True).order_by(Treinamento.titulo).all()
    if not treinamentos:
        flash(
            f'Nenhum treinamento ativo encontrado para o módulo {tipo_nr}.', 'warning')
    return render_template('treinamentos_por_modulo.html', treinamentos=treinamentos, modulo_nr=tipo_nr)


@app.route('/treinamentos/admin')
def treinamentos_admin():
    """Página administrativa para gerenciar treinamentos"""
    treinamentos = Treinamento.query.order_by(
        Treinamento.data_criacao.desc()).all()
    modelos_nr = ModeloNR.query.all()
    return render_template('treinamentos_admin.html',
                           treinamentos=treinamentos,
                           modelos_nr=modelos_nr)


@app.route('/treinamentos/upload', methods=['POST'])
def upload_treinamento():
    """Upload de novo vídeo de treinamento"""
    try:
        titulo = request.form['titulo'].strip()
        descricao = request.form.get('descricao', '').strip()
        tipo_nr = request.form['tipo_nr']
        duracao_minutos = request.form.get('duracao_minutos', type=int)
        obrigatorio = 'obrigatorio' in request.form

        if not titulo or not tipo_nr:
            flash('Título e Tipo de NR são obrigatórios!', 'error')
            return redirect(url_for('treinamentos_admin'))

        # Verificar se arquivo foi enviado
        if 'arquivo_video' not in request.files:
            flash('Nenhum arquivo de vídeo selecionado!', 'error')
            return redirect(url_for('treinamentos_admin'))

        arquivo = request.files['arquivo_video']
        if arquivo.filename == '':
            flash('Nenhum arquivo selecionado!', 'error')
            return redirect(url_for('treinamentos_admin'))

        # Validar extensão do arquivo
        extensoes_permitidas = {'.mp4', '.avi',
                                '.mov', '.wmv', '.mkv', '.webm'}
        if not any(arquivo.filename.lower().endswith(ext) for ext in extensoes_permitidas):
            flash(
                'Formato de vídeo não suportado! Use: MP4, AVI, MOV, WMV, MKV ou WEBM', 'error')
            return redirect(url_for('treinamentos_admin'))

        # Criar pasta de vídeos se não existir
        videos_dir = os.path.join('static', 'videos')
        os.makedirs(videos_dir, exist_ok=True)

        # Gerar nome único para o arquivo
        nome_arquivo = secure_filename(arquivo.filename)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        nome_arquivo = f"{timestamp}_{nome_arquivo}"
        caminho_arquivo = os.path.join(videos_dir, nome_arquivo)

        # Salvar arquivo
        arquivo.save(caminho_arquivo)

        # Criar registro no banco
        treinamento = Treinamento(
            titulo=titulo,
            descricao=descricao if descricao else None,
            tipo_nr=tipo_nr,
            duracao_minutos=duracao_minutos,
            arquivo_video=caminho_arquivo,
            obrigatorio=obrigatorio
        )

        db.session.add(treinamento)
        db.session.commit()

        flash(f'Treinamento "{titulo}" adicionado com sucesso!', 'success')
        return redirect(url_for('treinamentos_admin'))

    except Exception as e:
        flash(f'Erro ao fazer upload: {str(e)}', 'error')
        return redirect(url_for('treinamentos_admin'))


@app.route('/treinamentos/editar/<int:treinamento_id>', methods=['POST'])
def editar_treinamento(treinamento_id):
    """Editar treinamento existente"""
    treinamento = Treinamento.query.get_or_404(treinamento_id)

    try:
        treinamento.titulo = request.form['titulo'].strip()
        treinamento.descricao = request.form.get(
            'descricao', '').strip() or None
        treinamento.tipo_nr = request.form['tipo_nr']
        treinamento.duracao_minutos = request.form.get(
            'duracao_minutos', type=int)
        treinamento.obrigatorio = 'obrigatorio' in request.form
        treinamento.data_atualizacao = datetime.now()

        db.session.commit()
        flash(
            f'Treinamento "{treinamento.titulo}" atualizado com sucesso!', 'success')

    except Exception as e:
        flash(f'Erro ao atualizar: {str(e)}', 'error')

    return redirect(url_for('treinamentos_admin'))


@app.route('/treinamentos/excluir/<int:treinamento_id>', methods=['POST'])
def excluir_treinamento(treinamento_id):
    """Excluir/desativar treinamento"""
    treinamento = Treinamento.query.get_or_404(treinamento_id)

    try:
        # Verificar se há progresso de funcionários
        progresso_count = ProgressoTreinamento.query.filter_by(
            treinamento_id=treinamento_id).count()

        if progresso_count > 0:
            # Desativar ao invés de excluir se há progresso
            treinamento.ativo = False
            flash(
                f'Treinamento "{treinamento.titulo}" foi desativado (há progresso de funcionários).', 'warning')
        else:
            # Excluir arquivo e registro se não há progresso
            if os.path.exists(treinamento.arquivo_video):
                os.remove(treinamento.arquivo_video)
            db.session.delete(treinamento)
            flash(
                f'Treinamento "{treinamento.titulo}" excluído com sucesso!', 'success')

        db.session.commit()

    except Exception as e:
        flash(f'Erro ao excluir: {str(e)}', 'error')

    return redirect(url_for('treinamentos_admin'))


@app.route('/treinamentos/assistir/<int:treinamento_id>')
def assistir_treinamento(treinamento_id):
    """Página para assistir um treinamento específico"""
    treinamento = Treinamento.query.get_or_404(treinamento_id)

    # Por enquanto, vamos simular um funcionário logado (ID = 1)
    # Em produção, isso viria da sessão do usuário
    funcionario_id = 1  # TEMPORÁRIO - substituir por sessão real

    # Buscar ou criar progresso
    progresso = ProgressoTreinamento.query.filter_by(
        funcionario_id=funcionario_id,
        treinamento_id=treinamento_id
    ).first()

    if not progresso:
        progresso = ProgressoTreinamento(
            funcionario_id=funcionario_id,
            treinamento_id=treinamento_id
        )
        db.session.add(progresso)
        db.session.commit()

    # Buscar próximo treinamento da mesma NR
    proximo_treinamento = Treinamento.query.filter(
        Treinamento.tipo_nr == treinamento.tipo_nr,
        Treinamento.ativo == True,
        Treinamento.id != treinamento_id
    ).order_by(Treinamento.titulo).first()

    # Buscar outros treinamentos da mesma NR
    outros_treinamentos = Treinamento.query.filter(
        Treinamento.tipo_nr == treinamento.tipo_nr,
        Treinamento.ativo == True,
        Treinamento.id != treinamento_id
    ).order_by(Treinamento.titulo).limit(5).all()

    return render_template('assistir_treinamento.html',
                           treinamento=treinamento,
                           progresso=progresso,
                           proximo_treinamento=proximo_treinamento,
                           outros_treinamentos=outros_treinamentos)


@app.route('/treinamentos/progresso/<int:treinamento_id>', methods=['POST'])
def atualizar_progresso(treinamento_id):
    """API para atualizar progresso de visualização"""
    try:
        funcionario_id = 1  # TEMPORÁRIO - substituir por sessão real
        tempo_assistido = request.json.get('tempo_assistido', 0)
        progresso_percent = request.json.get('progresso_percent', 0)

        progresso = ProgressoTreinamento.query.filter_by(
            funcionario_id=funcionario_id,
            treinamento_id=treinamento_id
        ).first()

        if progresso:
            progresso.tempo_assistido_segundos = tempo_assistido
            progresso.progresso_percent = progresso_percent
            progresso.data_ultima_visualizacao = datetime.now()

            # Marcar como concluído se progresso >= 90%
            if progresso_percent >= 90 and not progresso.concluido:
                progresso.concluido = True
                progresso.data_conclusao = datetime.now()
                
                # Verificar se existe avaliação para este treinamento
                avaliacao = Avaliacao.query.filter_by(
                    treinamento_id=treinamento_id, 
                    ativo=True
                ).first()
                
                tem_avaliacao = avaliacao is not None
                
                # Verificar se já fez a avaliação
                ja_fez_avaliacao = False
                if avaliacao:
                    resultado_existente = ResultadoAvaliacao.query.filter_by(
                        funcionario_id=funcionario_id,
                        avaliacao_id=avaliacao.id
                    ).first()
                    ja_fez_avaliacao = resultado_existente is not None
                
                db.session.commit()
                
                return jsonify({
                    'status': 'success', 
                    'concluido': True,
                    'tem_avaliacao': tem_avaliacao,
                    'ja_fez_avaliacao': ja_fez_avaliacao,
                    'url_avaliacao': url_for('avaliacao_treinamento', treinamento_id=treinamento_id) if tem_avaliacao and not ja_fez_avaliacao else None
                })

            db.session.commit()

            return jsonify({'status': 'success', 'concluido': progresso.concluido})

        return jsonify({'status': 'error', 'message': 'Progresso não encontrado'})

    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)})


@app.route('/importar_funcionarios', methods=['POST'])
def importar_funcionarios():
    """Rota para importar funcionários em lote"""
    try:
        funcionarios_data = [
            {
                'nome': 'ANDRE ROBERTO DOS SANTOS',
                'cpf': '313.673.508-08',
                'data_nascimento': '12/01/1980',
                'telefone': '11 94870-7283',
                'email': 'andrealtatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Técnico de Segurança'
            },
            {
                'nome': 'BRUNO PINHEIRO DE LIMA',
                'cpf': '401.201.328-93',
                'data_nascimento': '27/04/1990',
                'telefone': '11 95273-9404',
                'email': 'brunoaltatelas@gmail.com',
                'senha': 'Alta972600#',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'CAIO VINICIUS PAES DO NASCIMENTO',
                'cpf': '414.038.588-02',
                'data_nascimento': '30/05/1991',
                'telefone': '11 95487-5236',
                'email': '',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'CARLOS ALBERTO DOS SANTOS',
                'cpf': '288.915.798-90',
                'data_nascimento': '21/11/1978',
                'telefone': '16 98865-1612',
                'email': 'carlosaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Supervisor de Instalação'
            },
            {
                'nome': 'CLEUTON CLEBER VIEIRA ROMANO',
                'cpf': '270.176.758-00',
                'data_nascimento': '05/12/1979',
                'telefone': '19 990058626',
                'email': 'cleutonaltatelas@gmail.com',
                'senha': 'Alta972600#',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'DIOGO PEREIRA CARDOSO',
                'cpf': '058.224.705-56',
                'data_nascimento': '30/12/1991',
                'telefone': '16 98103-7258',
                'email': 'diogoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'EDERSON GONÇALVES VIEIRA DA SILVA',
                'cpf': '495.685.328-97',
                'data_nascimento': '10/09/1999',
                'telefone': '16 99430-4340',
                'email': 'edersongonsalvesaltatelas@gmail.com',
                'senha': 'Alta972600',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'ERDESSON QUIRINO DE LIMA',
                'cpf': '385.283.058-30',
                'data_nascimento': '17/10/1988',
                'telefone': '16 98192-4980',
                'email': 'erdersonaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'EVERTON SAMPAIO MELO DA COSTA',
                'cpf': '258.334.658-00',
                'data_nascimento': '14/01/1976',
                'telefone': '11 96917-5480',
                'email': 'evertonaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Supervisor de Instalação'
            },
            {
                'nome': 'FELIPE ARAUJO DOS SANTOS',
                'cpf': '101.876.985-46',
                'data_nascimento': '13/03/2003',
                'telefone': '16 99292-3212',
                'email': 'felipealtatelas@gmail.com',
                'senha': 'Alta972600',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'FERNANDO SILVA DE SOUZA',
                'cpf': '005.378.972-59',
                'data_nascimento': '25/01/1989',
                'telefone': '16 99278-2701',
                'email': 'fernandoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'FILIPE DELFINO',
                'cpf': '350.969.688-39',
                'data_nascimento': '07/05/1987',
                'telefone': '16 99261-7738',
                'email': 'filipedelfinoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'GABRIEL JACKSON PEIXOTO RIBEIRO',
                'cpf': '438.515.568-28',
                'data_nascimento': '08/07/1996',
                'telefone': '16 99361-7195',
                'email': 'gabrielaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'GUSTAVO ADOLFO CASTANEDA CARDENAS',
                'cpf': '023.066.106-81',
                'data_nascimento': '23/12/1986',
                'telefone': '11 96293-9343',
                'email': 'gustavoaltatelas@gmail.com',
                'senha': 'Alta972600@',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'JARDIEL NUNES DA SILVA',
                'cpf': '070.368.533-36',
                'data_nascimento': '13/09/2000',
                'telefone': '16 98191-9551',
                'email': 'jardielaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'JUNIEL DE SOUSA',
                'cpf': '011.619.363-86',
                'data_nascimento': '07/09/1984',
                'telefone': '16 99204-7416',
                'email': 'junielaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'KAIQUE DE OLIVEIRA SANTOS',
                'cpf': '435.731.658-85',
                'data_nascimento': '02/12/1997',
                'telefone': '16 99447-8913',
                'email': 'kaiquealtatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'LEANDRO RODRIGUES DE SOUZA',
                'cpf': '414.455.308-64',
                'data_nascimento': '26/07/1987',
                'telefone': '16 99396-5158',
                'email': 'leandroaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'LEONARDO DOUGLAS DE OLIVEIRA CORREA',
                'cpf': '525.803.638-31',
                'data_nascimento': '05/03/2003',
                'telefone': '11 93397-5926',
                'email': 'leonardoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'LUCAS EDUARDO RIBEIRO FRANÇA E SILVA',
                'cpf': '468.944.148-07',
                'data_nascimento': '27/10/2006',
                'telefone': '16 98836-9753',
                'email': 'lucasaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Menor Aprendiz'
            },
            {
                'nome': 'LUCIANO DE JESUS',
                'cpf': '264.298.638-16',
                'data_nascimento': '06/11/1975',
                'telefone': '16 99754-1563',
                'email': 'lucianoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Supervisor de Instalação'
            },
            {
                'nome': 'LUIS HENRIQUE ROCETTI DA SILVA',
                'cpf': '412.241.338-97',
                'data_nascimento': '08/03/1992',
                'telefone': '16 99649-4762',
                'email': 'luisaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'MAICON DOUGLAS DOS SANTOS INFORZATO',
                'cpf': '463.539.228-76',
                'data_nascimento': '01/05/1993',
                'telefone': '16 98152-5822',
                'email': 'maiconaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'MATHEUS SOUSA COSTA',
                'cpf': '538.169.808-99',
                'data_nascimento': '01/12/2001',
                'telefone': '16 99111-6444',
                'email': 'matheusaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'MOEMIA COSTA DE OLIVEIRA',
                'cpf': '362.026.838-05',
                'data_nascimento': '05/02/1987',
                'telefone': '11 96150-8453',
                'email': 'moemia.costa@altatelas.com',
                'senha': 'Alta972600$',
                'funcao': 'Assistente Administrativo'
            },
            {
                'nome': 'PABLO HENRIQUE RIBEIRO FRIZI',
                'cpf': '584.518.788-57',
                'data_nascimento': '22/12/2007',
                'telefone': '16 99449-1102',
                'email': 'pabloaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Menor Aprendiz'
            },
            {
                'nome': 'PAULO FERREIRA CARDOSO',
                'cpf': '427.979.418-96',
                'data_nascimento': '27/10/1997',
                'telefone': '16 99715-6763',
                'email': 'pauloaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'RAFAEL DE SOUZA COSTA',
                'cpf': '106.957.083-41',
                'data_nascimento': '09/05/2005',
                'telefone': '89 9421-1070',
                'email': 'rafaelaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Menor Aprendiz'
            },
            {
                'nome': 'RODRIGO APARECIDO COELHO PAULISTA',
                'cpf': '218.937.248-83',
                'data_nascimento': '02/12/1982',
                'telefone': '16 99627-3207',
                'email': 'rodrigoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'RONALDO HENRIQUE DO CARMO',
                'cpf': '342.202.368-21',
                'data_nascimento': '14/09/1986',
                'telefone': '16 98846-2058',
                'email': 'ronaldoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'VICTOR HUGO VIANA GARCEZ',
                'cpf': '384.657.848-76',
                'data_nascimento': '15/09/2000',
                'telefone': '16 99307-5592',
                'email': 'victorhugoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'VITOR BRENO DE SOUZA ARAUJO',
                'cpf': '485.789.938-88',
                'data_nascimento': '15/09/2001',
                'telefone': '16 98187-3464',
                'email': 'vitorbrenoaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            },
            {
                'nome': 'WAGNER ZAMBONI',
                'cpf': '286.246.658-17',
                'data_nascimento': '17/08/1980',
                'telefone': '16 99722-6712',
                'email': 'wagneraltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Supervisor de Instalação'
            },
            {
                'nome': 'WENDEL HENRIQUE ANTONIO DA SILVA',
                'cpf': '504.069.548-95',
                'data_nascimento': '09/12/1998',
                'telefone': '16 99166-2527',
                'email': 'wendelaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Instalador de Telas'
            },
            {
                'nome': 'WILLIAM RODRIGUES DA SILVA',
                'cpf': '524.429.368-01',
                'data_nascimento': '16/03/2001',
                'telefone': '16 99650-8192',
                'email': 'williamaltatelas@gmail.com',
                'senha': 'Alta972600$',
                'funcao': 'Auxiliar de Instalação'
            }
        ]

        sucessos = 0
        erros = 0
        mensagens_erro = []

        for func_data in funcionarios_data:
            try:
                # Verificar se CPF já existe
                cpf_limpo = re.sub(r'[^0-9]', '', func_data['cpf'])
                cpf_formatado = f"{cpf_limpo[:3]}.{cpf_limpo[3:6]}.{cpf_limpo[6:9]}-{cpf_limpo[9:11]}"


                # Verificar se funcionário já existe e atualizar
                funcionario_existente = Funcionario.query.filter_by(
                    cpf=cpf_formatado).first()
                if funcionario_existente:
                    # Atualizar funcionário existente com nova função
                    funcionario_existente.funcao = func_data.get(
                        'funcao', 'Instalador de Telas')
                    funcionario_existente.telefone = func_data['telefone']
                    funcionario_existente.email = func_data['email'] if func_data['email'] else None
                    funcionario_existente.senha = func_data['senha']
                    sucessos += 1
                    continue

                # Converter data de nascimento
                data_nascimento = datetime.strptime(
                    func_data['data_nascimento'], '%d/%m/%Y').date()

                # Verificar se funcionário já existe e atualizar
                funcionario_existente = Funcionario.query.filter_by(
                    cpf=cpf_formatado).first()
                if funcionario_existente:
                    # Atualizar funcionário existente com nova função
                    funcionario_existente.funcao = func_data.get(
                        'funcao', 'Instalador de Telas')
                    funcionario_existente.telefone = func_data['telefone']
                    funcionario_existente.email = func_data['email'] if func_data['email'] else None
                    if func_data['senha']:
                        funcionario_existente.set_password(func_data['senha'])
                    sucessos += 1
                    continue

                # Criar funcionário
                funcionario = Funcionario(
                    nome=func_data['nome'],
                    cpf=cpf_formatado,
                    data_nascimento=data_nascimento,
                    telefone=func_data['telefone'],
                    email=func_data['email'] if func_data['email'] else None,
                    # Usar função específica ou padrão
                    funcao=func_data.get('funcao', 'Instalador de Telas'),
                    data_admissao=datetime.now().date()  # Data atual como admissão
                )
                
                # Define a senha
                if func_data['senha']:
                    funcionario.set_password(func_data['senha'])

                db.session.add(funcionario)
                sucessos += 1

            except Exception as e:
                erros += 1
                mensagens_erro.append(
                    f"Erro ao processar {func_data['nome']}: {str(e)}")

        db.session.commit()

        mensagem = f"Importação concluída! {sucessos} funcionários adicionados"
        if erros > 0:
            mensagem += f", {erros} erros encontrados"

        flash(mensagem, 'success' if erros == 0 else 'warning')

        if mensagens_erro:
            # Mostrar apenas os primeiros 5 erros
            for msg in mensagens_erro[:5]:
                flash(msg, 'error')

        return redirect(url_for('funcionarios'))

    except Exception as e:
        flash(f'Erro na importação: {str(e)}', 'error')
        return redirect(url_for('funcionarios'))

# Inicializar banco de dados


def criar_tabelas():
    db.create_all()

    # Criar modelos padrão se não existirem
    modelos_padrao = [
        {'tipo_nr': 'NR06', 'descricao': 'Equipamentos de Proteção Individual'},
        {'tipo_nr': 'NR10', 'descricao': 'Segurança em Instalações e Serviços em Eletricidade'},
        {'tipo_nr': 'NR12', 'descricao': 'Segurança no Trabalho em Máquinas e Equipamentos'},
        {'tipo_nr': 'NR18', 'descricao': 'Condições de Segurança na Indústria da Construção'},
        {'tipo_nr': 'NR33',
            'descricao': 'Segurança e Saúde nos Trabalhos em Espaços Confinados'},
        {'tipo_nr': 'NR35', 'descricao': 'Trabalho em Altura'}
    ]

    for modelo_info in modelos_padrao:
        if not ModeloNR.query.filter_by(tipo_nr=modelo_info['tipo_nr']).first():
            modelo = ModeloNR(
                tipo_nr=modelo_info['tipo_nr'],
                caminho_modelo_pptx=f"modelos_nr/{modelo_info['tipo_nr']}_modelo.pptx",
                descricao=modelo_info['descricao']
            )
            db.session.add(modelo)

    # Criar cargos padrão se não existirem
    cargos_padrao = [
        {'nome': 'Técnico de Segurança',
            'descricao': 'Responsável pela segurança do trabalho'},
        {'nome': 'Supervisor de Instalação',
            'descricao': 'Supervisiona equipes de instalação'},
        {'nome': 'Instalador de Telas',
            'descricao': 'Executa instalação de telas de proteção'},
        {'nome': 'Auxiliar de Instalação', 'descricao': 'Auxilia nas instalações'},
        {'nome': 'Menor Aprendiz', 'descricao': 'Aprendiz menor de idade'},
        {'nome': 'Assistente Administrativo',
            'descricao': 'Atividades administrativas'},
        {'nome': 'Gerente', 'descricao': 'Gerenciamento de equipes e projetos'},
        {'nome': 'Coordenador', 'descricao': 'Coordenação de atividades específicas'}
    ]

    for cargo_info in cargos_padrao:
        if not Cargo.query.filter_by(nome=cargo_info['nome'], ativo=True).first():
            cargo = Cargo(
                nome=cargo_info['nome'],
                descricao=cargo_info['descricao']
            )
            db.session.add(cargo)

    # Criar avaliações de exemplo se não existirem
    treinamento_nr06 = Treinamento.query.filter_by(tipo_nr='NR06').first()
    if treinamento_nr06:
        avaliacao_existente = Avaliacao.query.filter_by(treinamento_id=treinamento_nr06.id).first()
        if not avaliacao_existente:
            # Criar avaliação para NR06
            avaliacao_nr06 = Avaliacao(
                treinamento_id=treinamento_nr06.id,
                titulo='Avaliação NR06 - Equipamentos de Proteção Individual',
                descricao='Teste seus conhecimentos sobre EPIs e sua importância na segurança do trabalho.',
                nota_minima_aprovacao=7.0
            )
            db.session.add(avaliacao_nr06)
            db.session.flush()  # Para obter o ID
            
            # Perguntas da avaliação NR06
            perguntas_nr06 = [
                {
                    'texto': 'O que significa EPI?',
                    'ordem': 1,
                    'pontos': 2.0,
                    'opcoes': [
                        {'texto': 'Equipamento de Proteção Individual', 'correta': True, 'ordem': 1},
                        {'texto': 'Equipamento de Prevenção Individual', 'correta': False, 'ordem': 2},
                        {'texto': 'Equipamento de Proteção Integral', 'correta': False, 'ordem': 3},
                        {'texto': 'Equipamento de Prevenção Integral', 'correta': False, 'ordem': 4}
                    ]
                },
                {
                    'texto': 'Quando o EPI deve ser utilizado?',
                    'ordem': 2,
                    'pontos': 2.0,
                    'opcoes': [
                        {'texto': 'Apenas quando há risco iminente', 'correta': False, 'ordem': 1},
                        {'texto': 'Sempre que houver risco que não possa ser eliminado', 'correta': True, 'ordem': 2},
                        {'texto': 'Somente quando solicitado pelo supervisor', 'correta': False, 'ordem': 3},
                        {'texto': 'Apenas em caso de emergência', 'correta': False, 'ordem': 4}
                    ]
                },
                {
                    'texto': 'Quem é responsável por fornecer o EPI adequado?',
                    'ordem': 3,
                    'pontos': 2.0,
                    'opcoes': [
                        {'texto': 'O próprio trabalhador', 'correta': False, 'ordem': 1},
                        {'texto': 'O empregador', 'correta': True, 'ordem': 2},
                        {'texto': 'O sindicato', 'correta': False, 'ordem': 3},
                        {'texto': 'O governo', 'correta': False, 'ordem': 4}
                    ]
                },
                {
                    'texto': 'O que fazer se o EPI estiver danificado?',
                    'ordem': 4,
                    'pontos': 2.0,
                    'opcoes': [
                        {'texto': 'Continuar usando normalmente', 'correta': False, 'ordem': 1},
                        {'texto': 'Tentar consertá-lo', 'correta': False, 'ordem': 2},
                        {'texto': 'Comunicar imediatamente e solicitar substituição', 'correta': True, 'ordem': 3},
                        {'texto': 'Usar apenas em casos extremos', 'correta': False, 'ordem': 4}
                    ]
                },
                {
                    'texto': 'Qual a principal finalidade do capacete de segurança?',
                    'ordem': 5,
                    'pontos': 2.0,
                    'opcoes': [
                        {'texto': 'Proteger contra o sol', 'correta': False, 'ordem': 1},
                        {'texto': 'Proteger contra impactos na cabeça', 'correta': True, 'ordem': 2},
                        {'texto': 'Identificar o trabalhador', 'correta': False, 'ordem': 3},
                        {'texto': 'Melhorar a aparência', 'correta': False, 'ordem': 4}
                    ]
                }
            ]
            
            for pergunta_data in perguntas_nr06:
                pergunta = PerguntaAvaliacao(
                    avaliacao_id=avaliacao_nr06.id,
                    texto_pergunta=pergunta_data['texto'],
                    ordem=pergunta_data['ordem'],
                    pontos=pergunta_data['pontos']
                )
                db.session.add(pergunta)
                db.session.flush()
                
                for opcao_data in pergunta_data['opcoes']:
                    opcao = OpcaoResposta(
                        pergunta_id=pergunta.id,
                        texto_opcao=opcao_data['texto'],
                        correta=opcao_data['correta'],
                        ordem=opcao_data['ordem']
                    )
                    db.session.add(opcao)

    db.session.commit()


# ==================== ROTA PARA SERVIR VÍDEOS ====================

@app.route('/video/<path:filename>')
def serve_video(filename):
    """Servir vídeos com suporte a streaming e Range requests"""
    print(f"SERVE_VIDEO CHAMADA COM: {filename}")
    try:
        # Limpar o nome do arquivo para evitar problemas
        filename = secure_filename(filename)
        video_path = os.path.join('static', 'videos', filename)

        print(f"Tentando servir vídeo: {video_path}")
        print(f"Arquivo existe: {os.path.exists(video_path)}")

        if not os.path.exists(video_path):
            print(f"Arquivo não encontrado: {video_path}")
            raise NotFound()

        # Detectar tipo MIME
        mime_type, _ = mimetypes.guess_type(video_path)
        if not mime_type:
            mime_type = 'video/mp4'

        print(f"MIME type: {mime_type}")
        file_size = os.path.getsize(video_path)
        print(f"Tamanho do arquivo: {file_size} bytes")

        # Verificar se é uma requisição de Range
        range_header = request.headers.get('Range')
        print(f"Range header: {range_header}")

        if range_header:
            # Parse do header Range
            range_match = re.search(r'bytes=(\d+)-(\d*)', range_header)
            if range_match:
                start = int(range_match.group(1))
                end = int(range_match.group(2)) if range_match.group(
                    2) else file_size - 1

                # Garantir que end não ultrapasse o tamanho do arquivo
                end = min(end, file_size - 1)
                print(f"Servindo range: {start}-{end}")

                def generate():
                    with open(video_path, 'rb') as f:
                        f.seek(start)
                        remaining = end - start + 1

                        while remaining > 0:
                            chunk_size = min(8192, remaining)
                            chunk = f.read(chunk_size)
                            if not chunk:
                                break
                            yield chunk
                            remaining -= len(chunk)

                response = Response(
                    generate(),
                    206,  # Partial Content
                    headers={
                        'Content-Type': mime_type,
                        'Accept-Ranges': 'bytes',
                        'Content-Range': f'bytes {start}-{end}/{file_size}',
                        'Content-Length': str(end - start + 1),
                        'Cache-Control': 'no-cache'
                    }
                )
                return response

        # Resposta normal (arquivo completo)
        print("Servindo arquivo completo")
        return send_file(
            video_path,
            mimetype=mime_type,
            as_attachment=False,
            download_name=filename
        )

    except Exception as e:
        print(f"Erro ao servir vídeo {filename}: {str(e)}")
        import traceback
        traceback.print_exc()
        raise NotFound()


# ==================== ROTA PARA DEBUG DE VÍDEOS ====================

@app.route('/debug/videos')
def debug_videos():
    """Debug: listar todos os vídeos disponíveis"""
    videos_dir = os.path.join('static', 'videos')
    videos = []

    if os.path.exists(videos_dir):
        for filename in os.listdir(videos_dir):
            if filename.lower().endswith(('.mp4', '.avi', '.mov', '.wmv', '.mkv', '.webm')):
                filepath = os.path.join(videos_dir, filename)
                size = os.path.getsize(filepath)
                videos.append({
                    'filename': filename,
                    'size': f"{size / (1024*1024):.1f} MB",
                    'url': url_for('serve_video', filename=filename)
                })

    # Debug dos treinamentos no banco
    treinamentos = Treinamento.query.all()

    debug_info = {
        'videos_na_pasta': videos,
        'treinamentos_no_banco': [
            {
                'id': t.id,
                'titulo': t.titulo,
                'arquivo_video': t.arquivo_video,
                'existe': os.path.exists(t.arquivo_video) if t.arquivo_video else False
            } for t in treinamentos
        ]
    }

    return f"<pre>{str(debug_info)}</pre>"


@app.route('/teste-video')
def teste_video():
    """Página de teste para reprodução de vídeos"""
    return render_template('teste_video.html')


# ==================== SISTEMA DE AVALIAÇÃO ====================

@app.route('/treinamentos/<int:treinamento_id>/avaliacao')
def avaliacao_treinamento(treinamento_id):
    """Página de avaliação de um treinamento específico"""
    treinamento = Treinamento.query.get_or_404(treinamento_id)
    funcionario_id = 1  # TEMPORÁRIO - substituir por sessão real
    
    # Verificar se o funcionário completou o vídeo
    progresso = ProgressoTreinamento.query.filter_by(
        funcionario_id=funcionario_id,
        treinamento_id=treinamento_id
    ).first()
    
    if not progresso or not progresso.concluido:
        flash('Você precisa completar o treinamento antes de fazer a avaliação.', 'warning')
        return redirect(url_for('assistir_treinamento', treinamento_id=treinamento_id))
    
    # Buscar avaliação do treinamento
    avaliacao = Avaliacao.query.filter_by(treinamento_id=treinamento_id, ativo=True).first()
    
    if not avaliacao:
        flash('Nenhuma avaliação disponível para este treinamento.', 'info')
        return redirect(url_for('assistir_treinamento', treinamento_id=treinamento_id))
    
    # Verificar se já fez a avaliação
    resultado_existente = ResultadoAvaliacao.query.filter_by(
        funcionario_id=funcionario_id,
        avaliacao_id=avaliacao.id
    ).first()
    
    if resultado_existente:
        return redirect(url_for('resultado_avaliacao', resultado_id=resultado_existente.id))
    
    # Buscar perguntas ordenadas
    perguntas = PerguntaAvaliacao.query.filter_by(
        avaliacao_id=avaliacao.id
    ).order_by(PerguntaAvaliacao.ordem).all()
    
    return render_template('avaliacao_treinamento.html',
                         treinamento=treinamento,
                         avaliacao=avaliacao,
                         perguntas=perguntas)

@app.route('/treinamentos/<int:treinamento_id>/avaliacao/submeter', methods=['POST'])
def submeter_avaliacao(treinamento_id):
    """Submeter respostas da avaliação"""
    try:
        funcionario_id = 1  # TEMPORÁRIO - substituir por sessão real
        avaliacao = Avaliacao.query.filter_by(treinamento_id=treinamento_id, ativo=True).first_or_404()
        
        # Verificar se já fez a avaliação
        resultado_existente = ResultadoAvaliacao.query.filter_by(
            funcionario_id=funcionario_id,
            avaliacao_id=avaliacao.id
        ).first()
        
        if resultado_existente:
            flash('Você já realizou esta avaliação.', 'warning')
            return redirect(url_for('resultado_avaliacao', resultado_id=resultado_existente.id))
        
        # Processar respostas
        perguntas = PerguntaAvaliacao.query.filter_by(avaliacao_id=avaliacao.id).all()
        pontos_obtidos = 0
        pontos_totais = sum(p.pontos for p in perguntas)
        
        # Criar resultado
        resultado = ResultadoAvaliacao(
            funcionario_id=funcionario_id,
            avaliacao_id=avaliacao.id,
            nota_obtida=0,  # Será calculada
            aprovado=False,  # Será calculado
            tempo_realizacao_segundos=request.form.get('tempo_realizacao', type=int)
        )
        db.session.add(resultado)
        db.session.flush()  # Para obter o ID
        
        # Processar cada resposta
        for pergunta in perguntas:
            opcao_escolhida_id = request.form.get(f'pergunta_{pergunta.id}', type=int)
            
            if opcao_escolhida_id:
                opcao_escolhida = OpcaoResposta.query.get(opcao_escolhida_id)
                correta = opcao_escolhida.correta if opcao_escolhida else False

                if correta:
                    pontos_obtidos += pergunta.pontos
                
                # Salvar resposta
                resposta = RespostaFuncionario(
                    resultado_avaliacao_id=resultado.id,
                    pergunta_id=pergunta.id,
                    opcao_escolhida_id=opcao_escolhida_id,
                    correta=correta
                )
                db.session.add(resposta)
        
        # Calcular nota final (0-10)
        nota_final = (pontos_obtidos / pontos_totais) * 10 if pontos_totais > 0 else 0
        aprovado = nota_final >= avaliacao.nota_minima_aprovacao
        
        # Atualizar resultado
        resultado.nota_obtida = nota_final
        resultado.aprovado = aprovado
        
        db.session.commit()
        
        return redirect(url_for('resultado_avaliacao', resultado_id=resultado.id))
        
    except Exception as e:
        db.session.rollback()
        flash(f'Erro ao processar avaliação: {str(e)}', 'error')
        return redirect(url_for('avaliacao_treinamento', treinamento_id=treinamento_id))

@app.route('/avaliacao/resultado/<int:resultado_id>')
def resultado_avaliacao(resultado_id):
    """Página com resultado da avaliação"""
    resultado = ResultadoAvaliacao.query.get_or_404(resultado_id)
    
    # Verificar se o usuário pode ver este resultado
    funcionario_id = 1  # TEMPORÁRIO
    if resultado.funcionario_id != funcionario_id:
        flash('Acesso negado.', 'error')
        return redirect(url_for('index'))
    
    return render_template('resultado_avaliacao.html', resultado=resultado)

@app.route('/admin/avaliacoes')
def admin_avaliacoes():

    """Página administrativa para gerenciar avaliações"""
    avaliacoes = Avaliacao.query.order_by(Avaliacao.data_criacao.desc()).all()
    treinamentos = Treinamento.query.filter_by(ativo=True).all()
    
    return render_template('admin_avaliacoes.html', 
                         avaliacoes=avaliacoes,
                         treinamentos=treinamentos)

# ==================== SISTEMA DE AUTENTICAÇÃO ====================

@app.route('/login', methods=['GET', 'POST'])
def login():
    """Página de login"""
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        cpf = request.form['cpf'].replace('.', '').replace('-', '').strip()
        senha = request.form['senha']
        
        funcionario = Funcionario.query.filter_by(cpf=cpf).first()
        
        if funcionario and funcionario.check_password(senha):
            if not funcionario.ativo:
                flash('Sua conta está desativada. Entre em contato com o administrador.', 'error')
                return render_template('login.html')
            
            login_user(funcionario, remember=request.form.get('remember'))
            funcionario.data_ultimo_login = datetime.now()
            db.session.commit()
            
            # Redirecionar para próxima página ou dashboard
            next_page = request.args.get('next')
            if funcionario.primeiro_login:
                flash('Bem-vindo! Este é seu primeiro acesso. Considere alterar sua senha.', 'info')
                funcionario.primeiro_login = False
                db.session.commit()
                return redirect(next_page) if next_page else redirect(url_for('alterar_senha'))
            
            flash(f'Bem-vindo(a), {funcionario.nome}!', 'success')
            return redirect(next_page) if next_page else redirect(url_for('index'))
        else:
            flash('CPF ou senha incorretos.', 'error')
    
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    """Logout do usuário"""
    logout_user()
    flash('Você foi desconectado com sucesso.', 'info')
    return redirect(url_for('login'))

@app.route('/alterar-senha', methods=['GET', 'POST'])
@login_required
def alterar_senha():
    """Página para alterar senha"""
    if request.method == 'POST':
        senha_atual = request.form['senha_atual']
        nova_senha = request.form['nova_senha']
        confirmar_senha = request.form['confirmar_senha']
        
        if not current_user.check_password(senha_atual):
            flash('Senha atual incorreta.', 'error')
        elif nova_senha != confirmar_senha:
            flash('A confirmação da senha não confere.', 'error')
        elif len(nova_senha) < 6:
            flash('A nova senha deve ter pelo menos 6 caracteres.', 'error')
        else:
            current_user.set_password(nova_senha)
            current_user.primeiro_login = False
            db.session.commit()
            flash('Senha alterada com sucesso!', 'success')
            return redirect(url_for('index'))
    
    return render_template('alterar_senha.html')

@app.route('/perfil')
@login_required
def perfil():
    """Página do perfil do usuário"""
    return render_template('perfil.html', funcionario=current_user)

@app.route('/admin/criar-usuario', methods=['GET', 'POST'])
@login_required
def criar_usuario():
    """Página para administradores criarem novos usuários"""
    if not current_user.admin:
        flash('Acesso negado. Apenas administradores podem criar usuários.', 'error')
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        nome = request.form['nome'].strip()
        cpf = request.form['cpf'].replace('.', '').replace('-', '').strip()
        senha_inicial = request.form.get('senha_inicial', cpf[:6])  # Primeiros 6 dígitos do CPF como padrão
        admin = 'admin' in request.form
        
        # Verificar se CPF já existe
        if Funcionario.query.filter_by(cpf=cpf).first():
            flash('Já existe um funcionário com este CPF.', 'error')
        else:
            funcionario = Funcionario(
                nome=nome,
                cpf=cpf,
                admin=admin,
                funcao=request.form.get('funcao', 'Funcionário'),
                email=request.form.get('email', ''),
                telefone=request.form.get('telefone', '')
            )
            funcionario.set_password(senha_inicial)
            
            db.session.add(funcionario)
            db.session.commit()
            
            flash(f'Usuário {nome} criado com sucesso! Senha inicial: {senha_inicial}', 'success')
            return redirect(url_for('funcionarios'))
    
    return render_template('criar_usuario.html')

@app.route('/admin/reset-senha/<int:funcionario_id>', methods=['POST'])
@login_required
def reset_senha(funcionario_id):
    """Reset de senha por administrador"""
    if not current_user.admin:
        flash('Acesso negado.', 'error')
        return redirect(url_for('funcionarios'))
    
    funcionario = Funcionario.query.get_or_404(funcionario_id)
    nova_senha = funcionario.cpf[:6]  # Primeiros 6 dígitos do CPF
    
    funcionario.set_password(nova_senha)
    funcionario.primeiro_login = True
    db.session.commit()
    
    flash(f'Senha de {funcionario.nome} resetada para: {nova_senha}', 'success')
    return redirect(url_for('funcionarios'))


@login_manager.user_loader
def load_user(user_id):
    return Funcionario.query.get(int(user_id))

if __name__ == '__main__':
    with app.app_context():
        criar_tabelas()
    app.run(debug=True)
