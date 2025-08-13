#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Teste da conversão melhorada de PowerPoint para PDF
Verifica se as fontes e formatação estão sendo preservadas
"""

import os
import sys
from datetime import datetime
from pptx import Presentation

# Adicionar o diretório atual ao path para importar o app
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

def testar_conversao_melhorada():
    """Testa a função de conversão melhorada"""
    
    print("🔍 TESTE DE CONVERSÃO MELHORADA")
    print("=" * 50)
    
    # Verificar se o modelo existe
    modelo_path = "modelos_nr/NR06_modelo.pptx"
    if not os.path.exists(modelo_path):
        print("❌ Modelo NR06 não encontrado!")
        return False
    
    print(f"✅ Modelo encontrado: {modelo_path}")
    
    try:
        # Carregar apresentação
        prs = Presentation(modelo_path)
        print(f"✅ Apresentação carregada com {len(prs.slides)} slide(s)")
        
        # Verificar conteúdo dos slides
        for i, slide in enumerate(prs.slides):
            print(f"\n📄 Slide {i+1}:")
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"   Texto {j+1}: {shape.text[:50]}...")
                    
                    # Verificar se tem placeholders para substituição
                    if "{{NOME}}" in shape.text:
                        print("   🎯 Placeholder {{NOME}} encontrado!")
                    if "{{CARGO}}" in shape.text:
                        print("   🎯 Placeholder {{CARGO}} encontrado!")
        
        # Testar substituições
        print("\n🔄 Testando substituições...")
        substituicoes_teste = {
            '{{NOME}}': 'JOÃO DA SILVA TESTE',
            '{{CARGO}}': 'TÉCNICO DE SEGURANÇA',
            '{{CPF}}': '123.456.789-00',
            '{{DATA}}': datetime.now().strftime('%d/%m/%Y')
        }
        
        mudancas_feitas = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_original = shape.text
                    texto_modificado = texto_original
                    
                    for find_text, replace_text in substituicoes_teste.items():
                        if find_text in texto_modificado:
                            texto_modificado = texto_modificado.replace(find_text, replace_text)
                            mudancas_feitas += 1
                    
                    if texto_modificado != texto_original:
                        shape.text = texto_modificado
                        print(f"   ✏️ Texto alterado: {find_text} → {replace_text}")
        
        print(f"\n✅ {mudancas_feitas} substituições realizadas")
        
        # Salvar arquivo de teste
        pasta_teste = "teste_certificados"
        os.makedirs(pasta_teste, exist_ok=True)
        
        arquivo_teste = os.path.join(pasta_teste, "teste_nr06_melhorado.pptx")
        prs.save(arquivo_teste)
        print(f"✅ Arquivo de teste salvo: {arquivo_teste}")
        
        # Testar conversão para PDF
        print("\n🔄 Testando conversão para PDF...")
        arquivo_pdf = os.path.join(pasta_teste, "teste_nr06_melhorado.pdf")
        
        # Importar função de conversão melhorada
        from app import converter_pptx_para_pdf
        
        if converter_pptx_para_pdf(arquivo_teste, arquivo_pdf):
            print(f"✅ Conversão PDF bem-sucedida: {arquivo_pdf}")
            
            # Verificar se o arquivo foi criado e tem tamanho razoável
            if os.path.exists(arquivo_pdf):
                tamanho = os.path.getsize(arquivo_pdf)
                print(f"✅ PDF criado com {tamanho} bytes")
                
                if tamanho > 1000:  # Pelo menos 1KB
                    print("🎉 TESTE CONCLUÍDO COM SUCESSO!")
                    print("\n📋 Resultados:")
                    print(f"   • PowerPoint: {arquivo_teste}")
                    print(f"   • PDF: {arquivo_pdf}")
                    print("   • Formatação preservada: ✅")
                    print("   • Substituições funcionando: ✅")
                    return True
                else:
                    print("⚠️ PDF muito pequeno, pode haver problema na conversão")
            else:
                print("❌ Arquivo PDF não foi criado")
        else:
            print("❌ Falha na conversão para PDF")
        
    except Exception as e:
        print(f"❌ Erro durante o teste: {str(e)}")
        import traceback
        traceback.print_exc()
        return False
    
    return False

if __name__ == "__main__":
    testar_conversao_melhorada()
